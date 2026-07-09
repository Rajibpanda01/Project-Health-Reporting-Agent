from __future__ import annotations

import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

COLORS = {
    "canvas": RGBColor(255, 255, 255),
    "ink": RGBColor(17, 17, 17),
    "muted": RGBColor(95, 99, 104),
    "panel": RGBColor(241, 243, 244),
    "line": RGBColor(200, 205, 210),
    "accent": RGBColor(217, 72, 28),
    "accent_soft": RGBColor(253, 233, 226),
    "accent_warm": RGBColor(248, 215, 204),
}


def generate_monthly_presentation(deck_data_path: Path | str, output_pptx: Path | str) -> dict[str, str]:
    deck_data_path = Path(deck_data_path).resolve()
    output_pptx = Path(output_pptx).resolve()
    output_pptx.parent.mkdir(parents=True, exist_ok=True)

    with deck_data_path.open("r", encoding="utf-8") as handle:
        data = json.load(handle)

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    _build_executive_summary(_blank_slide(prs), data)
    _build_portfolio_health(_blank_slide(prs), data)
    _build_emerging_risks(_blank_slide(prs), data)
    _build_critical_projects(_blank_slide(prs), data)
    _build_recommendations(_blank_slide(prs))
    _build_outlook(_blank_slide(prs), data)

    prs.save(output_pptx)

    return {
        "pptx_path": str(output_pptx),
        "generator": "python-pptx",
    }


def _blank_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = COLORS["canvas"]
    return slide


def _add_text(
    slide,
    text: Any,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int = 18,
    color: RGBColor | None = None,
    bold: bool = False,
    align: PP_ALIGN | None = None,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)

    lines = str(text or "").split("\n")
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.name = "Aptos"
        paragraph.font.color.rgb = color or COLORS["ink"]
        if align is not None:
            paragraph.alignment = align
    return box


def _add_rect(slide, left: float, top: float, width: float, height: float, fill: RGBColor, line: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(0)
    return shape


def _add_rule(slide, left: float, top: float, width: float, fill: RGBColor | None = None, height: float = 0.02):
    _add_rect(slide, left, top, width, height, fill or COLORS["line"])


def _metric_card(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    value: Any,
    label: str,
    fill: RGBColor | None = None,
    value_color: RGBColor | None = None,
):
    _add_rect(slide, left, top, width, height, fill or COLORS["panel"])
    _add_text(slide, value, left + 0.18, top + 0.16, width - 0.36, 0.48, 34, value_color or COLORS["ink"], True)
    _add_text(slide, label, left + 0.18, top + 0.88, width - 0.36, height - 0.98, 16, COLORS["muted"])


def _add_title(slide, title: str, subtitle: str = ""):
    _add_text(slide, title, 0.5, 0.34, 9.1, 0.58, 34, COLORS["ink"], True)
    if subtitle:
        _add_text(slide, subtitle, 0.5, 0.96, 10.4, 0.45, 18, COLORS["muted"])


def _add_slide_number(slide, number: int):
    _add_text(slide, f"{number:02}", 12.32, 6.95, 0.5, 0.25, 12, COLORS["muted"], align=PP_ALIGN.RIGHT)


def _ordered_drivers(data: dict[str, Any]) -> list[tuple[str, int]]:
    drivers = data.get("overall", {}).get("lead_driver_counts", {})
    return sorted(drivers.items(), key=lambda item: item[1], reverse=True)


def _compact_project_name(name: str) -> str:
    return str(name or "").removeprefix("Zycus - ").replace(" Implementation", " implementation")


def _short_reason(text: str) -> str:
    return str(text or "").rstrip(".")


def _build_executive_summary(slide, data: dict[str, Any]) -> None:
    overall = data.get("overall", {})
    rag_counts = overall.get("rag_counts", {})
    drivers = _ordered_drivers(data)
    lead_driver = drivers[0][0] if drivers else "Schedule"
    second_driver = drivers[1][0] if len(drivers) > 1 else "Execution"
    focus_project = _compact_project_name(
        data.get("plan_b", {}).get("project_name")
        or (data.get("critical_projects") or [{}])[0].get("project_name")
        or "the highest-risk project"
    )

    _add_title(slide, "Executive Summary", "Automated monthly client-ready view generated from the analyzed project plans.")
    _metric_card(slide, 0.5, 1.62, 2.8, 1.56, overall.get("project_count", 0), "Projects")
    _metric_card(slide, 3.5, 1.62, 2.8, 1.56, rag_counts.get("Green", 0), "Green", COLORS["accent_soft"])
    _metric_card(slide, 6.5, 1.62, 2.8, 1.56, rag_counts.get("Amber", 0), "Amber")
    _metric_card(slide, 9.5, 1.62, 2.8, 1.56, rag_counts.get("Red", 0), "Red", COLORS["accent_warm"], COLORS["accent"])

    _add_rule(slide, 0.5, 3.62, 1.85, COLORS["ink"], 0.03)
    _add_text(
        slide,
        f"{overall.get('red_pct', 0)}% of analyzed views are currently red, which indicates that schedule recovery and execution control need executive attention before the next client reporting cycle.",
        0.5,
        3.88,
        7.9,
        1.2,
        22,
    )

    _add_rect(slide, 8.82, 3.68, 3.5, 3.23, COLORS["panel"])
    _add_text(slide, "This month at a glance", 9.08, 3.96, 2.95, 0.32, 20, bold=True)
    _add_text(
        slide,
        f"{lead_driver} is the most common lead driver.\n\n{second_driver} pressure is the next-most common concern across the current plans.\n\n{focus_project} remains red and should stay on the next escalation review list.",
        9.08,
        4.42,
        2.86,
        2.2,
        16,
    )
    _add_slide_number(slide, 1)


def _build_portfolio_health(slide, data: dict[str, Any]) -> None:
    overall = data.get("overall", {})
    portfolio = data.get("portfolio", {})
    rag_counts = overall.get("rag_counts", {})
    validation_accuracy = portfolio.get("validation_accuracy")
    validation_label = "N/A" if validation_accuracy is None else f"{round(validation_accuracy * 100)}%"
    validation_narrative = (
        "The original workbooks did not provide enough benchmark labels to calculate a validation accuracy figure."
        if validation_accuracy is None
        else f"Validation against the embedded workbook health labels holds at {round(validation_accuracy * 100)}%, which is directionally useful for leadership screening while keeping the reasoning transparent."
    )

    _add_title(
        slide,
        "Portfolio Health",
        "RAG distribution, average score, and model validation point to a portfolio that is broadly under delivery pressure.",
    )
    _add_rect(slide, 0.5, 1.6, 6.45, 4.94, COLORS["panel"])
    _add_text(slide, "RAG Distribution", 0.78, 1.84, 2.5, 0.3, 20, bold=True)
    _metric_card(slide, 0.96, 2.54, 1.62, 1.98, rag_counts.get("Green", 0), "Green")
    _metric_card(slide, 2.9, 2.54, 1.62, 1.98, rag_counts.get("Amber", 0), "Amber")
    _metric_card(slide, 4.84, 2.54, 1.62, 1.98, rag_counts.get("Red", 0), "Red", COLORS["accent_soft"], COLORS["accent"])
    _add_text(
        slide,
        f"{overall.get('project_count', 0)} project plans were analyzed in this run, and every current view is red or amber-weighted enough to need active leadership attention.",
        0.96,
        4.9,
        5.5,
        0.85,
        18,
    )

    _metric_card(slide, 7.46, 1.78, 2.42, 1.6, f"{overall.get('average_score', 0)}/100", "Overall Score")
    _metric_card(slide, 10.06, 1.78, 2.25, 1.6, validation_label, "Validation Accuracy", COLORS["accent_soft"], COLORS["accent"])
    _add_text(
        slide,
        f"The average score across all {overall.get('project_count', 0)} analyzed project plans is {overall.get('average_score', 0)}/100. {validation_narrative}",
        7.46,
        3.9,
        4.85,
        1.86,
        19,
    )
    _add_slide_number(slide, 2)


def _build_emerging_risks(slide, data: dict[str, Any]) -> None:
    portfolio = data.get("portfolio", {})
    _add_title(
        slide,
        "Emerging Risks",
        "Delay, budget burn, staffing pressure, and stakeholder confidence are emerging together rather than as isolated signals.",
    )
    _metric_card(slide, 0.5, 1.7, 2.8, 1.78, portfolio.get("burn_ahead_gt_10pts", 0), "Budget Burn")
    _metric_card(slide, 3.46, 1.7, 2.8, 1.78, portfolio.get("delay_gt_14_days", 0), "Schedule Delays", COLORS["accent_soft"], COLORS["accent"])
    _metric_card(slide, 0.5, 3.67, 2.8, 1.78, portfolio.get("resource_util_gt_90", 0), "Resource Issues")
    _metric_card(slide, 3.46, 3.67, 2.8, 1.78, portfolio.get("negative_or_neutral_sentiment", 0), "Stakeholder Concerns")

    _add_rect(slide, 6.7, 1.7, 5.6, 3.72, COLORS["panel"])
    _add_text(slide, "Risk interpretation", 7.02, 2.0, 2.4, 0.32, 22, bold=True)
    _add_text(
        slide,
        f"{portfolio.get('delay_gt_14_days', 0)} projects are already more than two weeks behind, and {portfolio.get('burn_ahead_gt_10pts', 0)} are spending ahead of delivery progress.\n\n{portfolio.get('critical_blocker_projects', 0)} projects also carry two or more critical blockers, which means the schedule risk is likely to stay elevated unless escalation gets faster.\n\nResource utilization above 90% and neutral or negative stakeholder sentiment show that the pressure is visible both inside the delivery teams and with client-facing sponsors.",
        7.02,
        2.48,
        4.75,
        2.45,
        17,
    )
    _add_slide_number(slide, 3)


def _build_critical_projects(slide, data: dict[str, Any]) -> None:
    critical_projects = data.get("critical_projects") or []
    _add_title(slide, "Critical Projects", "Top five red projects by score, ranked by the severity of the current delivery signal.")

    rows = max(len(critical_projects), 1) + 1
    table_shape = slide.shapes.add_table(rows, 3, Inches(0.5), Inches(1.72), Inches(11.82), Inches(4.38))
    table = table_shape.table
    table.columns[0].width = Inches(3.35)
    table.columns[1].width = Inches(1.15)
    table.columns[2].width = Inches(7.32)

    headers = ["Project", "Score", "Reason"]
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        _format_cell(cell, COLORS["panel"], 13, True)

    if critical_projects:
        for row, item in enumerate(critical_projects, start=1):
            values = [item.get("project_name", ""), item.get("overall_score", ""), _short_reason(item.get("reason_1", ""))]
            for col, value in enumerate(values):
                cell = table.cell(row, col)
                cell.text = str(value)
                _format_cell(cell, COLORS["canvas"] if row % 2 else COLORS["panel"], 11, False)
    else:
        cell = table.cell(1, 0)
        cell.text = "No critical projects identified"
        _format_cell(cell, COLORS["canvas"], 11, False)
        for col in (1, 2):
            _format_cell(table.cell(1, col), COLORS["canvas"], 11, False)

    _add_text(
        slide,
        "These projects should be reviewed first in the next escalation cadence because they combine the highest risk scores with clear delivery explanations.",
        0.5,
        6.32,
        10.3,
        0.42,
        18,
    )
    _add_slide_number(slide, 4)


def _format_cell(cell, fill: RGBColor, font_size: int, bold: bool) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = fill
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.05)
    cell.margin_bottom = Inches(0.05)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = COLORS["ink"]


def _build_recommendations(slide) -> None:
    _add_title(
        slide,
        "Recommendations",
        "The next operating cycle should focus on recovery discipline, blocker escalation, and cleaner execution reporting.",
    )
    actions = [
        (
            "01",
            "Increase staffing where utilization is already above safe limits",
            "Reassign or add capacity to the most constrained plans before schedule recovery work pushes the same teams further behind.",
        ),
        (
            "02",
            "Escalate blocked projects through a daily unblock cadence",
            "Use a standing review for critical blockers and P1 issues so execution stalls do not sit inside weekly status reports.",
        ),
        (
            "03",
            "Rebaseline delayed milestones and enforce dated recovery plans",
            "Reset dates at milestone level for the top red projects and review burn-versus-progress weekly until the gap narrows.",
        ),
    ]

    for index, (number, title, body) in enumerate(actions):
        top = 1.7 + index * 1.6
        fill = COLORS["accent_soft"] if index == 0 else COLORS["panel"]
        _add_rect(slide, 0.5, top, 11.82, 1.32, fill)
        _add_text(slide, number, 0.8, top + 0.38, 0.65, 0.34, 22, COLORS["accent"], True)
        _add_text(slide, title, 1.68, top + 0.25, 9.6, 0.36, 20, bold=True)
        _add_text(slide, body, 1.68, top + 0.66, 9.5, 0.44, 16)
    _add_slide_number(slide, 5)


def _build_outlook(slide, data: dict[str, Any]) -> None:
    outlook = data.get("outlook", {})
    plan_b = data.get("plan_b", {})
    improvements = list(outlook.get("expected_improvements") or [])
    while len(improvements) < 3:
        improvements.append("Keep leadership review focused on the highest-risk delivery signals.")

    _add_title(
        slide,
        "Next Month Outlook",
        "The near-term forecast improves only if the current red watchlist receives faster recovery actions than it did this month.",
    )
    _add_rect(slide, 0.5, 1.75, 5.12, 3.46, COLORS["panel"])
    _add_text(slide, "Forecast", 0.78, 2.06, 1.9, 0.3, 22, bold=True)
    _add_text(slide, outlook.get("forecast", ""), 0.78, 2.56, 4.35, 1.95, 20)

    _add_rect(slide, 6.02, 1.75, 6.3, 3.46, COLORS["accent_soft"])
    _add_text(slide, "Expected improvements", 6.32, 2.06, 4.0, 0.34, 20, bold=True)
    _add_text(
        slide,
        f"1. {improvements[0]}\n\n2. {improvements[1]}\n\n3. {improvements[2]}",
        6.32,
        2.64,
        5.4,
        2.13,
        16,
    )

    _add_text(
        slide,
        f"Project Plan B should stay on the watchlist next month because it remains {plan_b.get('rag_status', 'under review')} and still carries a {round(plan_b.get('progress_signal_gap') or 0)}-point gap between weekly and task-level progress.",
        0.5,
        5.82,
        11.82,
        0.58,
        19,
    )
    _add_slide_number(slide, 6)
